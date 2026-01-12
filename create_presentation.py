from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(26, 58, 82)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
TEXT_DARK = RGBColor(26, 26, 26)

def add_blue_title_slide(title, subtitle):
    """Add title slide with blue background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_white_slide(title, bullets):
    """Add white slide with title and bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)

def add_blue_slide(title, bullets):
    """Add blue slide with title and bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

def add_two_col_slide(title, left_title, left_bullets, right_title, right_bullets):
    """Add two column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    for bullet in left_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    for bullet in right_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

# SLIDE 1: Title
add_blue_title_slide("SciCon", "Scientific Conference Management Platform")

# SLIDE 2: Challenge
add_two_col_slide(
    "The Challenge",
    "Current State",
    [
        "Manual reviewer assignment",
        "Inconsistent paper evaluation",
        "Scattered registrations",
        "Time-consuming certificates",
        "Limited engagement"
    ],
    "SciCon Solution",
    [
        "Automated management",
        "Standardized scoring",
        "Centralized platform",
        "Auto-generated certs",
        "Interactive features"
    ]
)

# SLIDE 3: Tech Stack
add_two_col_slide(
    "Technology Stack",
    "Frontend",
    ["React 18", "Vite", "Shadcn/UI", "Tailwind CSS", "Axios"],
    "Backend",
    ["Django 4.2", "DRF", "JWT Auth", "PostgreSQL", "Docker"]
)

# SLIDE 4: Architecture
add_white_slide(
    "System Architecture",
    [
        "React SPA → JWT Token → Django REST API → ORM → PostgreSQL",
        "",
        "Frontend: Role-based routing, Shadcn/UI, centralized API",
        "",
        "Backend: 50+ REST endpoints, 14 models, custom permissions"
    ]
)

# SLIDE 5: User Roles
add_blue_slide(
    "Four User Roles",
    [
        "👨‍💼 Organizer: Create events, assign reviewers, manage registrations",
        "",
        "📝 Author: Submit papers, track reviews, get feedback",
        "",
        "🔍 Reviewer: Evaluate papers, score submissions",
        "",
        "🎫 Participant: Browse events, register, download certificates"
    ]
)

# SLIDE 6: Event Lifecycle
add_white_slide(
    "Event Lifecycle",
    [
        "Draft → Open Call → Reviewing → Program Ready → Ongoing → Completed",
        "",
        "Open Call: Authors submit papers with abstracts",
        "",
        "Reviewing: Organizer assigns reviewers for evaluation",
        "",
        "Program Ready: Accepted papers assigned to sessions",
        "",
        "Completed: Certificates generated automatically"
    ]
)

# SLIDE 7: Review Process
add_white_slide(
    "Intelligent Review Process",
    [
        "Three-Dimensional Scoring (1-5 scale):",
        "",
        "• Relevance: Alignment with conference themes",
        "",
        "• Quality: Scientific rigor and methodology",
        "",
        "• Originality: Contribution to the field",
        "",
        "Each reviewer scores independently"
    ]
)

# SLIDE 8: Decision Logic
add_white_slide(
    "Automated Decision Making",
    [
        "After 2+ reviews, system calculates average score:",
        "",
        "if avg_score >= 4.0  → ACCEPTED ✓",
        "if avg_score < 2.5   → REJECTED ✗",
        "otherwise            → REVISION REQUESTED 🔄",
        "",
        "Benefits: No bias • Consistency • Saves time"
    ]
)

# SLIDE 9: Database
add_blue_slide(
    "Core Models",
    [
        "14 interconnected database models:",
        "",
        "User • Event • Session • Submission • Review",
        "Registration • Workshop • Question • Survey",
        "Certificate • Message • Notification",
        "",
        "Event → Submissions → Reviews",
        "Event → Registrations → Users"
    ]
)

# SLIDE 10: Frontend
add_white_slide(
    "Frontend Organization",
    [
        "Feature-based modular structure:",
        "",
        "features/auth/ - Login, registration",
        "features/dashboard/ - 4 role dashboards",
        "features/events/ - Event management",
        "features/submissions/ - Paper workflow",
        "components/ui/ - Shadcn/UI atoms",
        "services/api.js - API client"
    ]
)

# SLIDE 11: Components
add_white_slide(
    "Shadcn/UI Components",
    [
        "Card: Dashboard panels, event listings",
        "",
        "Button: Actions with variants",
        "",
        "Badge: Status indicators",
        "",
        "Tabs: Dashboard navigation",
        "",
        "Toast: Notifications",
        "",
        "Table: Submission lists"
    ]
)

# SLIDE 12: API
add_white_slide(
    "REST API Structure",
    [
        "Over 50 endpoints:",
        "",
        "Auth: register, login, refresh, profile",
        "Events: list, create, edit, delete, stats",
        "Submissions: submit, list, assign-reviewers",
        "Reviews: evaluate, score, decide",
        "Registrations: register, payment, tracking",
        "Certificates: generate, download"
    ]
)

# SLIDE 13: Registration
add_blue_slide(
    "Registration Workflow",
    [
        "User clicks 'Register for Event'",
        "↓",
        "POST /api/events/:id/registrations/",
        "↓",
        "Backend validates & checks for duplicates",
        "↓",
        "Toast notification with success/error",
        "",
        "Complete in ~500ms with instant feedback"
    ]
)

# SLIDE 14: Security
add_white_slide(
    "Security Architecture",
    [
        "🔐 Authentication: JWT tokens, PBKDF2 hashing",
        "",
        "🛡️ Authorization: Role-Based Access Control",
        "",
        "✅ Data Integrity: Unique constraints, validation",
        "",
        "📁 File Security: Upload isolation, verification",
        "",
        "All sensitive operations protected by custom permissions"
    ]
)

# SLIDE 15: Certificates
add_white_slide(
    "Automatic Certificate Generation",
    [
        "System generates 4 certificate types:",
        "",
        "• Participation: All registered participants",
        "",
        "• Presentation: Accepted paper authors",
        "",
        "• Scientific Committee: Reviewers",
        "",
        "• Organization: Event organizers",
        "",
        "Generated automatically after event completion"
    ]
)

# SLIDE 16: Analytics
add_white_slide(
    "Event Analytics",
    [
        "Submissions by country/institution",
        "",
        "Registration metrics: participants, speakers",
        "",
        "Survey results: ratings, feedback",
        "",
        "Geographic distribution of attendees",
        "",
        "Session attendance tracking",
        "",
        "Real-time organizer dashboard"
    ]
)

# SLIDE 17: Development
add_two_col_slide(
    "Development Setup",
    "Tools",
    ["Git/GitHub", "Trello/Notion", "Discord/Slack"],
    "Team Roles",
    ["Project Manager", "Frontend Lead", "Backend Lead", "Database Admin"]
)

# SLIDE 18: Getting Started
add_white_slide(
    "Getting Started",
    [
        "Backend:",
        "python -m venv venv",
        "pip install -r requirements.txt",
        "python manage.py runserver",
        "",
        "Frontend:",
        "npm install && npm run dev"
    ]
)

# SLIDE 19: Timeline
add_white_slide(
    "Project Timeline",
    [
        "Phase 1: Requirements (Weeks 1-3)",
        "Phase 2: Backend API (Weeks 4-7)",
        "Phase 3: Frontend (Weeks 8-11)",
        "Phase 4: Integration & Testing (Weeks 12-13)",
        "Phase 5: Documentation & Deployment (Weeks 14-15)",
        "",
        "Duration: January 2025 → June 2026"
    ]
)

# SLIDE 20: Deliverables
add_blue_slide(
    "Project Deliverables",
    [
        "Complete source code (frontend + backend)",
        "",
        "Technical documentation with diagrams",
        "",
        "API endpoint documentation",
        "",
        "Database schema diagrams",
        "",
        "Deployment instructions",
        "",
        "GitHub repository with README"
    ]
)

# SLIDE 21: Conclusion
add_white_slide(
    "Key Takeaways",
    [
        "Enterprise-grade conference management platform",
        "",
        "Modern React + Django architecture",
        "",
        "Intelligent peer review with automation",
        "",
        "4 distinct dashboards for different roles",
        "",
        "Scalable from workshops to international conferences"
    ]
)

# Save
try:
    prs.save('SciCon_Presentation.pptx')
    print("✅ SUCCESS! Created: SciCon_Presentation.pptx")
    print("📂 Check your current directory for the file")
except Exception as e:
    print(f"❌ Error: {e}")